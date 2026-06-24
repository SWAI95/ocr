"""OCR 엔진 비교 발표자료(PPTX) 생성 — 동우 화인켐 브랜드.

data/measurements.json(성능) + data/pptx_assets/*.png(시각화) + APEX_Introduce/assets
(로고) 소비 → OCR_엔진_비교_발표.pptx (≤10장). 한글=맑은 고딕, 액센트 red/blue.

실행: ./.venv/bin/python scripts/build_pptx.py
"""
from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "data" / "pptx_assets"
LOGO_DIR = ROOT.parent / "APEX_Introduce" / "assets"
MEAS = ROOT / "data" / "measurements.json"
OUT = ROOT / "OCR_엔진_비교_발표.pptx"

# ── 브랜드 ────────────────────────────────────────────────────────────────
RED = RGBColor(0xE0, 0x30, 0x30)
RED7 = RGBColor(0xB6, 0x24, 0x24)
REDW = RGBColor(0xFC, 0xE5, 0xE5)
BLUE = RGBColor(0x00, 0x70, 0xB0)
BLUE7 = RGBColor(0x00, 0x56, 0x88)
BLUEW = RGBColor(0xE0, 0xF0, 0xF8)
INK = RGBColor(0x10, 0x10, 0x18)
G900 = RGBColor(0x18, 0x18, 0x1B)
G700 = RGBColor(0x3F, 0x3F, 0x46)
G500 = RGBColor(0x71, 0x71, 0x7A)
G200 = RGBColor(0xE4, 0xE4, 0xE7)
G50 = RGBColor(0xFA, 0xFA, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "맑은 고딕"

EMU = 914400
SW, SH = int(13.333 * EMU), int(7.5 * EMU)


def load_meas() -> dict:
    return json.loads(MEAS.read_text()) if MEAS.exists() else {}


M = load_meas()
# fallback(부분 데이터/미측정 대비) — 이번 세션 GPU 실측값
GPU_FB = {
    "hybrid_vl": {"doc_00": 1.52, "doc_04": 0.71, "doc_05": 1.56},
    "paddle": {"doc_00": 2.16, "doc_04": 1.16, "doc_05": 2.45},
    "paddle_vl": {"doc_00": 7.46, "doc_04": 0.72, "doc_05": 2.03},
    "ollama_vl": {"doc_00": 2.80, "doc_04": 1.89, "doc_05": 1.53},
    "easyocr": {"doc_00": 11.4, "doc_04": 8.02, "doc_05": 8.89},
    "tesseract": {"doc_00": 33.1, "doc_04": 28.7, "doc_05": 16.4},
}


def gpu_cer(eng, doc):
    try:
        return M["gpu"][doc][eng]["cer"]
    except Exception:
        return GPU_FB.get(eng, {}).get(doc)


def gpu_wer(eng, doc):
    try:
        return M["gpu"][doc][eng]["wer"]
    except Exception:
        return None


def fmt(v, suf="%"):
    return f"{v:.1f}{suf}" if isinstance(v, (int, float)) else "—"


# ── 저수준 헬퍼 ───────────────────────────────────────────────────────────
def set_font(run, name=FONT, size=18, bold=False, color=INK, italic=False):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", name)


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    return s


def rect(slide, x, y, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)),
                                Emu(int(w)), Emu(int(h)))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line:
        sp.line.color.rgb = line; sp.line.width = Pt(0.75)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def tbox(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    """lines: [(text,size,bold,color), ...] each becomes a paragraph."""
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, item in enumerate(lines):
        text, size, bold, color = (item + (INK,))[:4] if len(item) < 4 else item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)
        r = p.add_run(); r.text = text
        set_font(r, size=size, bold=bold, color=color)
    return tb


def eyebrow(slide, text, x=Emu(int(0.7 * EMU)), y=Emu(int(0.45 * EMU))):
    tb = slide.shapes.add_textbox(x, y, Emu(int(8 * EMU)), Emu(int(0.4 * EMU)))
    r = tb.text_frame.paragraphs[0].add_run(); r.text = text.upper()
    set_font(r, size=12, bold=True, color=G500)


def heading(slide, text, sub=None):
    eyebrow(slide, "DONGWOO FINE-CHEM · OCR")
    tb = slide.shapes.add_textbox(Emu(int(0.7 * EMU)), Emu(int(0.78 * EMU)),
                                  Emu(int(12 * EMU)), Emu(int(1.0 * EMU)))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = text
    set_font(r, size=30, bold=True, color=INK)
    rect(slide, 0.72 * EMU, 1.62 * EMU, 0.55 * EMU, 0.07 * EMU, RED)
    if sub:
        tb2 = slide.shapes.add_textbox(Emu(int(0.7 * EMU)), Emu(int(1.72 * EMU)),
                                       Emu(int(12 * EMU)), Emu(int(0.5 * EMU)))
        r2 = tb2.text_frame.paragraphs[0].add_run(); r2.text = sub
        set_font(r2, size=14, color=G700)
    # 코너 로고 마크
    mark = LOGO_DIR / "logo-mark.png"
    if mark.exists():
        slide.shapes.add_picture(str(mark), Emu(int(12.35 * EMU)), Emu(int(0.4 * EMU)),
                                 height=Emu(int(0.45 * EMU)))


def table(slide, x, y, w, headers, rows, col_w=None, best_row=None,
          fs=12, header_bg=INK):
    nrow, ncol = len(rows) + 1, len(headers)
    h = Emu(int((0.34 + 0.32 * len(rows)) * EMU))
    gtbl = slide.shapes.add_table(nrow, ncol, Emu(int(x)), Emu(int(y)),
                                  Emu(int(w)), h).table
    if col_w:
        for j, cw in enumerate(col_w):
            gtbl.columns[j].width = Emu(int(cw))
    for j, htxt in enumerate(headers):
        c = gtbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = header_bg
        c.margin_top = Pt(2); c.margin_bottom = Pt(2)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = htxt; set_font(r, size=fs, bold=True, color=WHITE)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = gtbl.cell(i + 1, j)
            c.fill.solid()
            c.fill.fore_color.rgb = REDW if best_row == i else (G50 if i % 2 else WHITE)
            c.margin_top = Pt(1); c.margin_bottom = Pt(1)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(val)
            set_font(r, size=fs, bold=(best_row == i),
                     color=INK if best_row == i else G700)
    return gtbl


def bullets(slide, x, y, w, items, fs=14, gap=6):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(4 * EMU)))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        if isinstance(it, tuple):
            lead, rest = it
            r1 = p.add_run(); r1.text = "● "; set_font(r1, size=fs, bold=True, color=RED)
            r2 = p.add_run(); r2.text = lead; set_font(r2, size=fs, bold=True, color=INK)
            r3 = p.add_run(); r3.text = rest; set_font(r3, size=fs, color=G700)
        else:
            r1 = p.add_run(); r1.text = "● "; set_font(r1, size=fs, bold=True, color=BLUE)
            r2 = p.add_run(); r2.text = it; set_font(r2, size=fs, color=G700)
    return tb


def card(slide, x, y, w, h, title, body, accent=BLUE):
    rect(slide, x, y, w, h, G50, line=G200)
    rect(slide, x, y, 0.08 * EMU, h, accent)
    tbox(slide, x + 0.2 * EMU, y + 0.14 * EMU, w - 0.35 * EMU, h - 0.2 * EMU,
         [(title, 15, True, INK)] + [(b, 12, False, G700) for b in body])


# ════════════════════════════════════════════════════════════════════════
def s_title(prs):
    s = blank(prs)
    logo = LOGO_DIR / "logo-dwfc.png"
    if logo.exists():
        slide_pic = s.shapes.add_picture(str(logo), Emu(int(0.7 * EMU)), Emu(int(0.7 * EMU)),
                                         height=Emu(int(0.55 * EMU)))
    rect(s, 0.72 * EMU, 3.05 * EMU, 0.7 * EMU, 0.09 * EMU, RED)
    tbox(s, 0.7 * EMU, 3.25 * EMU, 11.5 * EMU, 2.2 * EMU, [
        ("오프라인 한국어 문서 OCR", 40, True, INK),
        ("엔진·모델 비교 결과 보고", 40, True, BLUE),
    ])
    tbox(s, 0.72 * EMU, 5.35 * EMU, 11 * EMU, 1.2 * EMU, [
        ("테스트셋: 계약서 보관 신청서 3종 (케이아이에스 · 유테크솔루션 · 리스너)", 14, False, G700),
        ("측정 환경: RTX 5090 GPU  ·  지표: CER / WER  ·  완전 오프라인(내부망)", 13, False, G500),
    ])
    mark = LOGO_DIR / "logo-mark.png"
    if mark.exists():
        s.shapes.add_picture(str(mark), Emu(int(10.5 * EMU)), Emu(int(5.3 * EMU)),
                             height=Emu(int(1.4 * EMU)))


def s_dl_vs_vlm(prs):
    s = blank(prs)
    heading(s, "두 갈래의 OCR: 딥러닝 방식 vs VLM 방식",
            "OCR = 이미지 속 글자를 컴퓨터가 읽어 텍스트로 바꾸는 기술")
    y = 2.35 * EMU
    card(s, 0.7 * EMU, y, 5.9 * EMU, 2.5 * EMU, "딥러닝 방식 OCR (예: PaddleOCR)", [
        "글자 위치를 찾고(검출) → 한 줄씩 읽는(인식) 전용 신경망.",
        "매우 빠름 — 페이지당 0.3초 수준, 완전 오프라인.",
        "표·폼 구조에 강함. 같은 입력엔 항상 같은 결과(결정론적).",
        "약점: 줄바꿈·읽기 순서가 흔들릴 수 있음.",
    ], accent=BLUE)
    card(s, 6.85 * EMU, y, 5.8 * EMU, 2.5 * EMU, "VLM 방식 (비전언어모델, 예: Qwen2.5-VL)", [
        "이미지를 통째로 '보고 이해해' 글자를 읽는 AI.",
        "문장·읽기 순서가 자연스럽고 정확도가 높음.",
        "약점: 느림 — 페이지당 ~10초. 가끔 글자를 지어내기도.",
        "헤더/로고 같은 영역을 놓치기도 함.",
    ], accent=RED)
    rect(s, 0.7 * EMU, 5.25 * EMU, 11.95 * EMU, 1.35 * EMU, BLUEW)
    tbox(s, 0.95 * EMU, 5.4 * EMU, 11.5 * EMU, 1.1 * EMU, [
        ("핵심 차이 — 속도 ↔ 정확도의 맞교환", 16, True, BLUE7),
        ("빠른 응답이 중요하면 딥러닝 방식, 정확도가 중요하면 VLM 방식. "
         "Paddle은 바이두의 오픈소스 OCR로 검출·인식·표·레이아웃을 한 번에 처리한다.", 13, False, G700),
    ])


def s_three_stages(prs):
    s = blank(prs)
    heading(s, "인식의 3단계: 검출 · 인식 · 레이아웃",
            "OCR은 한 번에 읽는 게 아니라 단계로 나눠 처리한다")
    cards = [
        ("① 검출 (Detection)", ["글자가 \"어디\" 있는지", "위치(상자)를 찾는다."], BLUE),
        ("② 인식 (Recognition)", ["그 상자 안이 \"무슨\"", "글자인지 읽는다."], RED),
        ("③ 레이아웃 (Layout)", ["페이지를 본문·표·제목·", "머리글·그림·도장 구역으로 나눈다."], BLUE7),
    ]
    for i, (t, b, c) in enumerate(cards):
        card(s, (0.7 + i * 2.15) * EMU, 2.3 * EMU, 2.0 * EMU, 1.6 * EMU, t, b, accent=c)
    img = ASSETS / "region_overlay.png"
    if img.exists():
        s.shapes.add_picture(str(img), Emu(int(7.5 * EMU)), Emu(int(2.05 * EMU)),
                             height=Emu(int(4.7 * EMU)))
    tbox(s, 0.7 * EMU, 4.2 * EMU, 6.4 * EMU, 2.4 * EMU, [
        ("→ 레이아웃 덕분에 \"여기는 표, 여기는 제목\"을 구분", 14, True, INK),
        ("오른쪽은 실제 인식 화면 — 색 상자가 구역 종류를 표시한다.", 13, False, G700),
        ("파랑=본문, 빨강=제목, 회색=머리글, 진한 파랑=표.", 12, False, G500),
        ("표·그림·도장을 따로 떼어내 처리할 수 있어, 단순히 글자만 읽는 것보다 "
         "문서 구조를 살린 결과를 얻는다.", 13, False, G700),
    ])


def s_engines(prs):
    s = blank(prs)
    heading(s, "사용한 엔진 6종 한눈에 보기")
    rows = [
        ["hybrid_vl ★", "Paddle + VLM 병합", "폼·표=Paddle / 줄글=VLM. 장점만 결합 — 종합 1위", "딥러닝+VLM"],
        ["paddle", "PP-StructureV3", "딥러닝. 빠르고 표 구조에 강함. 완전 오프라인", "딥러닝"],
        ["paddle_vl", "PaddleOCR-VL 0.9B", "단독 VLM(서버 불필요). 읽기순서 강점", "VLM"],
        ["ollama_vl", "Qwen2.5-VL (로컬)", "로컬 Ollama VLM. 줄글 가독성 좋음", "VLM"],
        ["easyocr", "CRAFT + CRNN", "가벼운 전통 OCR. 문서 정확도는 낮음", "딥러닝"],
        ["tesseract", "Tesseract 5", "오픈소스 기본 비교군. 한국어 정확도 낮음", "딥러닝"],
    ]
    table(s, 0.7 * EMU, 2.3 * EMU, 11.95 * EMU,
          ["엔진", "구성", "특징", "방식"], rows,
          col_w=[1.9 * EMU, 2.6 * EMU, 6.0 * EMU, 1.45 * EMU], best_row=0, fs=12)
    tbox(s, 0.7 * EMU, 6.35 * EMU, 12 * EMU, 0.6 * EMU, [
        ("★ hybrid_vl = 자사 기본 엔진. 'Paddle의 정확한 글자' + 'VLM의 자연스러운 "
         "문장 배치'를 문자 단위로 합쳐 정확도·가독성을 동시에 잡았다.", 12, False, G700)])


def s_samples_perf(prs):
    s = blank(prs)
    heading(s, "샘플 문서 3종 & 모델별 정확도",
            "doc_00=케이아이에스 · doc_04=유테크솔루션 · doc_05=리스너 (모두 계약서 보관 신청서)")
    rect(s, 0.7 * EMU, 2.25 * EMU, 5.9 * EMU, 0.95 * EMU, BLUEW)
    tbox(s, 0.9 * EMU, 2.34 * EMU, 5.6 * EMU, 0.8 * EMU, [
        ("CER = 글자 오류율", 13, True, BLUE7),
        ("100글자 중 몇 자가 틀렸나(공백 무시). 낮을수록 좋음.", 11, False, G700)])
    rect(s, 6.75 * EMU, 2.25 * EMU, 5.9 * EMU, 0.95 * EMU, REDW)
    tbox(s, 6.95 * EMU, 2.34 * EMU, 5.6 * EMU, 0.8 * EMU, [
        ("WER = 단어(띄어쓰기) 오류율", 13, True, RED7),
        ("문장 가독성 지표. 한국어는 보조 지표로 본다.", 11, False, G700)])
    engs = ["hybrid_vl", "paddle", "paddle_vl", "ollama_vl", "easyocr", "tesseract"]
    rows = []
    for e in engs:
        rows.append([e + (" ★" if e == "hybrid_vl" else ""),
                     fmt(gpu_cer(e, "doc_00")), fmt(gpu_cer(e, "doc_04")),
                     fmt(gpu_cer(e, "doc_05"))])
    table(s, 0.7 * EMU, 3.45 * EMU, 8.2 * EMU,
          ["엔진", "doc_00", "doc_04", "doc_05"], rows,
          col_w=[2.5 * EMU, 1.9 * EMU, 1.9 * EMU, 1.9 * EMU], best_row=0, fs=12)
    tbox(s, 9.1 * EMU, 3.5 * EMU, 3.5 * EMU, 3 * EMU, [
        ("읽는 법 (CER, 낮을수록 좋음)", 13, True, INK),
        ("· hybrid_vl이 3개 문서 모두 최저 → 종합 1위", 12, False, G700),
        ("· 깨끗한 문서(doc_04)는 1% 미만으로 매우 정확", 12, False, G700),
        ("· tesseract는 한국어 문서에 부적합(참고용)", 12, False, G500),
        ("· 측정: RTX 5090 GPU, 문서 전체 페이지 기준", 11, False, G500),
    ])


def s_errors(prs):
    s = blank(prs)
    heading(s, "주요 오인식 케이스", "왜 틀리는지 알면 어디에 쓸 수 있는지 보인다")
    items = [
        ("원화기호 ₩ → 라틴 'W'  ", "도장처럼 ₩의 가로줄을 못 읽어 'W100,000'으로 오인식. "
         "→ '숫자 앞 W'를 자동으로 ₩로 교정(규칙 후처리)."),
        ("상단 헤더·로고 누락  ", "VLM이 페이지 맨 위 회사 로고·번호를 통째로 놓치는 경우. "
         "→ hybrid는 Paddle이 헤더를 잡아 보완."),
        ("표·폼 읽기순서 뒤섞임  ", "빈 칸 많은 폼에서 셀 순서가 엉켜 글자가 맞아도 순서가 틀림. "
         "→ 페이지를 통째로 읽는 VLM이 복원."),
        ("도장이 글자를 가림  ", "빨간 도장이 글자 위에 찍혀 그 아래 글자가 깨짐. "
         "→ 빨강 억제 전처리로 복구(다음 장)."),
    ]
    bullets(s, 0.7 * EMU, 2.4 * EMU, 12 * EMU, items, fs=15, gap=12)


def s_gpu_cpu(prs):
    s = blank(prs)
    heading(s, "GPU vs CPU — 같은 문서, 다른 속도",
            "doc_00(5페이지) 처리 시간. 정확도는 동일, 속도만 차이")
    engs = ["paddle", "easyocr", "tesseract", "ollama_vl", "paddle_vl", "hybrid_vl"]
    rows = []
    for e in engs:
        g = (M.get("gpu", {}).get("doc_00", {}).get(e) or {})
        c = (M.get("cpu_doc00", {}).get(e) or {})
        gs = f"{g['per_page_s']}s" if g.get("per_page_s") is not None else "—"
        cs = f"{c['per_page_s']}s" if c.get("per_page_s") is not None else "측정중"
        spd = "—"
        if g.get("per_page_s") and c.get("per_page_s"):
            spd = f"{c['per_page_s']/max(g['per_page_s'],0.01):.0f}배"
        rows.append([e, gs, cs, spd])
    table(s, 0.7 * EMU, 2.4 * EMU, 7.6 * EMU,
          ["엔진", "GPU (초/page)", "CPU (초/page)", "CPU 둔화"], rows,
          col_w=[2.2 * EMU, 1.9 * EMU, 1.9 * EMU, 1.6 * EMU], fs=12)
    rect(s, 8.5 * EMU, 2.4 * EMU, 4.15 * EMU, 3.7 * EMU, G50, line=G200)
    tbox(s, 8.7 * EMU, 2.6 * EMU, 3.8 * EMU, 3.4 * EMU, [
        ("환경이 속도를 좌우한다", 15, True, INK),
        ("· 딥러닝 엔진(paddle)은 CPU에서도 실용 가능", 12, False, G700),
        ("· VLM은 CPU에서 페이지당 수십~수백 초 → 사실상 비실용", 12, False, RED7),
        ("· 빠른 현장 응답이 필요하면 GPU가 사실상 필수", 12, False, G700),
        ("· 결론: 서비스 요구(속도)에 맞춰 엔진·하드웨어를 함께 정해야 함", 12, False, G700),
    ])


def s_seal(prs):
    s = blank(prs)
    heading(s, "도장 억제 — 양날의 검", "빨간 도장을 지우면 좋아지기도, 나빠지기도 한다")
    b = ASSETS / "doc04_stamp_before.png"; a = ASSETS / "doc04_stamp_after.png"
    if b.exists():
        s.shapes.add_picture(str(b), Emu(int(0.7 * EMU)), Emu(int(2.4 * EMU)), height=Emu(int(2.7 * EMU)))
        tbox(s, 0.7 * EMU, 5.15 * EMU, 2.0 * EMU, 0.4 * EMU, [("억제 전(도장이 글자 가림)", 10, False, G500)])
    if a.exists():
        s.shapes.add_picture(str(a), Emu(int(2.75 * EMU)), Emu(int(2.4 * EMU)), height=Emu(int(2.7 * EMU)))
        tbox(s, 2.75 * EMU, 5.15 * EMU, 2.0 * EMU, 0.4 * EMU, [("억제 후(빨강 제거→글자 복구)", 10, False, G500)])
    seal = M.get("seal", {})

    def srow(doc):
        out = []
        for e in ("hybrid_vl", "paddle", "paddle_vl"):
            d = seal.get(doc, {}).get(e, {})
            off, on = d.get("off_cer"), d.get("on_cer")
            out.append([e, fmt(off), fmt(on),
                        (f"{on-off:+.1f}%p" if off is not None and on is not None else "측정중")])
        return out
    tbox(s, 5.0 * EMU, 2.25 * EMU, 7.5 * EMU, 0.4 * EMU,
         [("doc_04 (도장이 글자 가림) — 억제하면 개선 ↓", 13, True, BLUE7)])
    table(s, 5.0 * EMU, 2.65 * EMU, 7.5 * EMU, ["엔진", "억제 전", "억제 후", "변화"],
          srow("doc_04"), col_w=[2.4 * EMU, 1.7 * EMU, 1.7 * EMU, 1.7 * EMU], fs=11)
    tbox(s, 5.0 * EMU, 4.35 * EMU, 7.5 * EMU, 0.4 * EMU,
         [("doc_00 (도장이 방해 안 함) — 억제하면 오히려 저하 ↑", 13, True, RED7)])
    table(s, 5.0 * EMU, 4.75 * EMU, 7.5 * EMU, ["엔진", "억제 전", "억제 후", "변화"],
          srow("doc_00"), col_w=[2.4 * EMU, 1.7 * EMU, 1.7 * EMU, 1.7 * EMU], fs=11)
    rect(s, 0.7 * EMU, 6.25 * EMU, 11.95 * EMU, 0.6 * EMU, REDW)
    tbox(s, 0.95 * EMU, 6.33 * EMU, 11.5 * EMU, 0.45 * EMU, [
        ("→ 전역 적용은 금물. 도장이 글자를 가리는 문서에만 선택적으로 적용해야 한다.", 13, True, RED7)])


def s_usecases(prs):
    s = blank(prs)
    heading(s, "이 기술로 할 수 있는 것 — 활용 방안",
            "인식된 문서 데이터를 어디에 쓰나")
    left = [
        ("문서 검색·질의(RAG) ", "인식 텍스트를 Vector DB에 넣어, 수천 건 계약서를 '검색'하고 "
         "LLM에게 \"이 조항 있는 계약 찾아줘\"라고 질문."),
        ("현장 즉시 인식 ", "작업자가 현장에서 문서를 촬영하면 즉시 텍스트화 → 그 데이터를 바로 업무에 활용."),
        ("문서 이상 탐지 ", "필수 항목 누락·금액 불일치·서명 누락 등을 자동 점검."),
        ("LLM 연동 자동화 ", "추출 정보로 요약·분류·번역·보고서 자동 생성."),
    ]
    right = [
        ("MSDS·CAS 추출 ", "물질안전보건자료에서 화학물질명·CAS 번호·위험문구 자동 수집."),
        ("시험성적서(COA) 자동화 ", "원료·제품 시험성적서 수치를 인식해 품질 시스템에 자동 입력."),
        ("규제·인증 문서 ", "인허가·환경·안전 규제 문서를 텍스트화해 관리·검색."),
        ("작업표준서·설비 문서 ", "현장 SOP·도면 주기재를 디지털화해 검색·교육에 활용."),
    ]
    tbox(s, 0.7 * EMU, 2.3 * EMU, 5.9 * EMU, 0.4 * EMU, [("공통 활용", 15, True, BLUE7)])
    bullets(s, 0.7 * EMU, 2.75 * EMU, 5.9 * EMU, left, fs=13, gap=10)
    tbox(s, 6.85 * EMU, 2.3 * EMU, 5.8 * EMU, 0.4 * EMU, [("화학 제조업 특화", 15, True, RED7)])
    bullets(s, 6.85 * EMU, 2.75 * EMU, 5.8 * EMU, right, fs=13, gap=10)


def s_limits(prs):
    s = blank(prs)
    heading(s, "한계점 & 도입을 위한 하드웨어 가이드")
    bullets(s, 0.7 * EMU, 2.25 * EMU, 6.2 * EMU, [
        ("VLM은 느리다 ", "정확하지만 GPU로도 페이지당 ~10초. 속도↔정확도를 골라야 함."),
        ("1글자 오류도 치명적인 분야엔 아직 부족 ", "완전 무인 처리보다 '사람 확인 보조'로 시작 권장."),
        ("추가 학습 필요 ", "서비스를 확정한 뒤, 해당 문서 샘플로 추가 학습해 인식률을 끌어올려야 함."),
        ("환경 편차 ", "CPU/GPU에 따라 속도가 수십 배 차이. 도입 전 환경 산정 필수."),
    ], fs=13, gap=11)
    tbox(s, 7.1 * EMU, 2.05 * EMU, 5.5 * EMU, 0.4 * EMU,
         [("이용 규모별 권장 사양 & 대략 견적", 14, True, INK)])
    rows = [
        ["소규모/현장 1~수명\n(딥러닝 위주)", "RTX 4060 Ti 16GB", "32GB / 8코어", "150~250만"],
        ["중규모/팀\n(VLM 가끔)", "RTX 4090 24GB\n또는 5090 32GB", "64GB / 12코어", "400~700만"],
        ["대규모/상시 서버\n(VLM 상시)", "RTX 6000 Ada 48GB\n/ L40S (다중 GPU)", "128GB+ / 16코어+", "1,500만~"],
        ["CPU 전용\n(딥러닝, 저처리량)", "GPU 없음", "32GB / 16코어", "~100만"],
    ]
    table(s, 7.1 * EMU, 2.5 * EMU, 5.55 * EMU,
          ["구분", "GPU", "RAM / CPU", "견적(원)"], rows,
          col_w=[1.8 * EMU, 1.85 * EMU, 1.35 * EMU, 0.55 * EMU], fs=9.5)
    tbox(s, 7.1 * EMU, 5.9 * EMU, 5.55 * EMU, 1.0 * EMU, [
        ("· 현재 RTX 5090 기준: 딥러닝 0.3초/page, VLM ~10초/page.", 11, False, G700),
        ("· VLM 상시 서비스는 GPU VRAM 24GB+ 권장. 견적은 시점·환율 따라 변동(대략치).", 10, False, G500)])


def main():
    prs = Presentation()
    prs.slide_width = Emu(SW); prs.slide_height = Emu(SH)
    s_title(prs)
    s_dl_vs_vlm(prs)
    s_three_stages(prs)
    s_engines(prs)
    s_samples_perf(prs)
    s_errors(prs)
    s_gpu_cpu(prs)
    s_seal(prs)
    s_usecases(prs)
    s_limits(prs)
    prs.save(str(OUT))
    print(f"[OK] {OUT}  슬라이드 {len(prs.slides.__iter__.__self__._sldIdLst)}장")


if __name__ == "__main__":
    main()
